import streamlit as st  
import requests  
from bs4 import BeautifulSoup  
import google.generativeai as genai  
from pptx import Presentation  
from pptx.util import Inches  
from fpdf import FPDF  
import os  
  
# Configure Google Generative AI  
genai.configure(api_key="AIzaSyCvM4yzyrUflRJdug-E9wtV_0ALWCwVGY0")  
model = genai.GenerativeModel('gemini-1.5-flash')  
  
# Memory to store citations  
agent_memory = {}  
  
# PDF Generator Class  
class PDFExporter:  
    def __init__(self):  
        self.pdf = FPDF()  
        self.pdf.set_auto_page_break(auto=True, margin=15)  
        self.pdf.add_page()  
        self.pdf.set_font("Arial", size=12)  
  
    def add_abstract(self, title, abstract, url):  
        # Handle potential Unicode issues by sanitizing text  
        title = self.sanitize_text(title)  
        abstract = self.sanitize_text(abstract)  
        url = self.sanitize_text(url)  
          
        self.pdf.set_font("Arial", 'B', 12)  
        self.pdf.cell(200, 10, txt=title, ln=True)  
        self.pdf.set_font("Arial", size=11)  
        self.pdf.multi_cell(0, 10, f"{abstract}\nLink: {url}\n\n")  
  
    def export(self, filename):  
        self.pdf.output(filename)  
  
    def sanitize_text(self, text):  
        """Sanitize text to remove unsupported characters."""  
        return text.encode('latin1', 'replace').decode('latin1')  
  
# Scraper Agent  
def fetch_pubmed_articles(disease, limit=3):  
    query = f"reducing {disease} readmission"  
    search_url = f"https://pubmed.ncbi.nlm.nih.gov/?term={query.replace(' ', '+')}"  
    response = requests.get(search_url)  
    soup = BeautifulSoup(response.text, "html.parser")  
    articles = soup.select(".docsum-content")[:limit]  
  
    results = []  
    for article in articles:  
        title_tag = article.select_one("a.docsum-title")  
        if not title_tag:  
            continue  
        title = title_tag.get_text(strip=True)  
        article_url = "https://pubmed.ncbi.nlm.nih.gov" + title_tag["href"]  
        abstract = fetch_abstract(article_url)  
        results.append({"title": title, "url": article_url, "abstract": abstract})  
    return results  
  
def fetch_abstract(url):  
    response = requests.get(url)  
    soup = BeautifulSoup(response.text, "html.parser")  
    abstract_tag = soup.select_one(".abstract-content.selected")  
    return abstract_tag.get_text(strip=True) if abstract_tag else "No abstract available."  
  
# Summarizer Agent  
def summarize_with_gemini(disease, articles):  
    full_text = "\n\n".join([f"Title: {a['title']}\nAbstract: {a['abstract']}" for a in articles])  
    prompt = f"""You are an expert summarizer. Your task is to extract key clinical and administrative insights.  
Disease: {disease}  
Summarize the following abstracts:  
{full_text}  
Output two sections:  
1. Clinical Strategy: (for doctors)  
2. Administrative Actions: (for hospital administrators)"""  
    response = model.generate_content(prompt)  
    return response.text.strip()  
  
# Deck Builder  
def create_deck(summaries):  
    prs = Presentation()  
    for disease, content in summaries.items():  
        slide_layout = prs.slide_layouts[1]  
        slide = prs.slides.add_slide(slide_layout)  
        slide.shapes.title.text = disease  
        slide.placeholders[1].text = content  
  
    # Add a final slide summarizing all diseases  
    final_slide = prs.slides.add_slide(prs.slide_layouts[1])  
    final_slide.shapes.title.text = "Comprehensive Readmission Program"  
    final_slide.placeholders[1].text = "\n\n".join([f"{d}: Summary Available" for d in summaries])  
  
    # Save the PowerPoint presentation  
    prs.save("readmission_summary_deck.pptx")  
  
# ====== STREAMLIT UI ======  
st.title("📚 Agentic AI - Medical Readmission Strategy Generator")  
st.markdown("Searches PubMed, summarizes using Gemini, and generates PDF + PPTX outputs.")  
  
# Input field for diseases  
diseases = st.text_input("Enter diseases (comma separated)", "CHF, Sepsis, UTI, Kidney failure")  
  
# Button to trigger the generation process  
if st.button("Generate Summary & Deck"):  
    st.info("Processing...")  
    disease_list = [d.strip() for d in diseases.split(",")]  
    summaries = {}  
    pdf = PDFExporter()  
  
    # Process each disease  
    for disease in disease_list:  
        st.write(f"🔍 Processing: **{disease}**")  
        articles = fetch_pubmed_articles(disease)  
        summary = summarize_with_gemini(disease, articles)  
        summaries[disease] = summary  
  
        # Store citations in memory  
        agent_memory[disease] = [(a['title'], a['url']) for a in articles]  
  
        # Add articles to PDF  
        for a in articles:  
            pdf.add_abstract(a['title'], a['abstract'], a['url'])  
  
        # Display the summary in the Streamlit UI  
        st.markdown(f"**{disease} Summary:**")  
        st.markdown(summary)  
  
    # Save outputs  
    create_deck(summaries)  
    pdf.export("pubmed_abstracts.pdf")  
  
    # Provide download buttons for the generated files  
    with open("pubmed_abstracts.pdf", "rb") as f:  
        st.download_button("📄 Download Abstracts PDF", f, file_name="pubmed_abstracts.pdf")  
    with open("readmission_summary_deck.pptx", "rb") as f:  
        st.download_button("📊 Download PowerPoint Deck", f, file_name="readmission_summary_deck.pptx")  
  
    st.success("✅ All done!")  
  
    # Display memory (citations)  
    st.markdown("### 🔗 Citations (Memory)")  
    for disease, citations in agent_memory.items():  
        st.markdown(f"**{disease}**")  
        for title, link in citations:  
            st.markdown(f"- [{title}]({link})")  
